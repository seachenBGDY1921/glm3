# encoding:utf-8
import argparse
import json
import logging
import os
import shutil
from datetime import datetime
from typing import List, Optional
import urllib
import asyncio
import nltk
import pydantic
import requests
import uvicorn
from apscheduler.triggers.interval import IntervalTrigger
from fastapi import Body, Request, FastAPI, File, Form, Query, UploadFile, WebSocket
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pydantic.typing import Dict
from typing_extensions import Annotated
from starlette.responses import RedirectResponse
from fastapi.responses import FileResponse

from chains.local_doc_qa import LocalDocQA
from configs.model_config import (KB_ROOT_PATH, EMBEDDING_DEVICE,
                                  EMBEDDING_MODEL, NLTK_DATA_PATH,
                                  VECTOR_SEARCH_TOP_K, LLM_HISTORY_LEN, OPEN_CROSS_DOMAIN)
import models.shared as shared
from models.loader.args import parser
from models.loader import LoaderCheckPoint

import os
import fitz  # PyMuPDF库，用于处理PDF
from openpyxl import load_workbook  # 用于处理Excel
from pptx import Presentation  # 用于处理PPT
import docx  # 用于处理Word


# nltk.data.path = [NLTK_DATA_PATH] + nltk.data.path

import redis
import jwt
import pymysql

redisdb = redis.Redis( host="dm.redis.hubpd-internal.com", port = 6379, password = 'gX92t@Yv3Xo')
host = "dm.hubpd.com"
#redisdb = redis.Redis(host='10.19.1.1', port=7001, password='myredis')
from dbutils.pooled_db import PooledDB


from datetime import datetime

#创建数据库连接池
pool = PooledDB(
    creator=pymysql,  # 数据库连接库
    maxconnections=20,  # 连接池最大连接数
    mincached=5,  # 初始化时，连接池中至少创建的空闲的连接
    maxcached=5,  # 连接池中最多共享的连接数量
    maxshared=3,  # 连接池中最多共享的连接数量
    blocking=True,  # 连接池中如果没有可用共享连接后，是否阻塞等待
    maxusage=None,  # 一个连接最多被重复使用的次数，None 表示无限制
    setsession=[],  # 开始会话前执行的命令列表
    ping=0,
    host='pdmi-db-1.rwlb.rds.aliyuncs.com',
    user='llm',
    password='EbcLvd2=YLL*',
    database='ai_dataset',
    charset='utf8',
    port=3306,
)


class BaseResponse(BaseModel):
    code: int = pydantic.Field(200, description="HTTP status code")
    msg: str = pydantic.Field("success", description="HTTP status message")

    class Config:
        schema_extra = {
            "example": {
                "code": 200,
                "msg": "success",
            }
        }

class UIDResponse(BaseResponse):
    last_id: int = pydantic.Field(0, description="HTTP status message")
    class Config:
        schema_extra = {
            "example": {
                "code": 200,
                "msg": "success",
                "last_id": 100
            }
        }



class UserDetailResponse(BaseResponse):
    user: Dict[str, str] = pydantic.Field("user:detail", description="HTTP status message")

    class Config:
        schema_extra = {
            "example": {
                "code": 200,
                "msg": "success",
                "user": {
                    "user_id": 1,
                    "user_name": "testname",
                    "password": "123",
                    "knowledge_base_id": "0",
                    "perms": "1",
                    "phone_num": "123456778",
                    "create_time": "2023-08-16 13:14:33",
                }
            }
        }



class ListDocsResponse(BaseResponse):
    data: List[str] = pydantic.Field(..., description="List of document names")

    class Config:
        schema_extra = {
            "example": {
                "code": 200,
                "msg": "success",
                "data": ["doc1.docx", "doc2.pdf", "doc3.txt"],
            }
        }


class ListDocsCreateResponse(BaseResponse):
    data: List[dict] = pydantic.Field(..., description="List of document names")
    count: int = pydantic.Field(0, description="知识库文件数量")

    class Config:
        schema_extra = {
            "example": {
                "code": 200,
                "msg": "success",
                "data": [{"doc1.docx": "2023-09-09"}, {"doc2.pdf": "2023-09-09"}, {"doc3.txt": "2023-09-09"}],
            }
        }


class ChatMessage(BaseModel):
    question: str = pydantic.Field(..., description="Question text")
    response: str = pydantic.Field(..., description="Response text")
    history: List[List[Optional[str]]] = pydantic.Field(..., description="History text")
    source_documents: List[dict] = pydantic.Field(
        ..., description="List of source documents and their scores"
    )

    class Config:
        schema_extra = {
            "example": {
                "question": "工伤保险如何办理？",
                "response": "根据已知信息，可以总结如下：\n\n1. 参保单位为员工缴纳工伤保险费，以保障员工在发生工伤时能够获得相应的待遇。\n2. 不同地区的工伤保险缴费规定可能有所不同，需要向当地社保部门咨询以了解具体的缴费标准和规定。\n3. 工伤从业人员及其近亲属需要申请工伤认定，确认享受的待遇资格，并按时缴纳工伤保险费。\n4. 工伤保险待遇包括工伤医疗、康复、辅助器具配置费用、伤残待遇、工亡待遇、一次性工亡补助金等。\n5. 工伤保险待遇领取资格认证包括长期待遇领取人员认证和一次性待遇领取人员认证。\n6. 工伤保险基金支付的待遇项目包括工伤医疗待遇、康复待遇、辅助器具配置费用、一次性工亡补助金、丧葬补助金等。",
                "history": [
                    [
                        "工伤保险是什么？",
                        "工伤保险是指用人单位按照国家规定，为本单位的职工和用人单位的其他人员，缴纳工伤保险费，由保险机构按照国家规定的标准，给予工伤保险待遇的社会保险制度。",
                    ]
                ],
                "source_documents": [
                    "出处 [1] 广州市单位从业的特定人员参加工伤保险办事指引.docx：\n\n\t( 一)  从业单位  (组织)  按“自愿参保”原则，  为未建 立劳动关系的特定从业人员单项参加工伤保险 、缴纳工伤保 险费。",
                    "出处 [2] ...",
                    "出处 [3] ...",
                ],
            }
        }


class DBMessage(BaseModel):
    total_count: int = pydantic.Field(..., description="Question text")
    total_pages: int = pydantic.Field(..., description="Response text")
    dbresult: list[dict[str,str]] = pydantic.Field(..., description="History text")
    count: Optional[int] = pydantic.Field(None, description="自构建数据集数量")
    class Config:
        schema_extra = {
            "example": {
                "total_count": 1,
                "total_pages": 2,
                "history": [(1, 'tag1', 'description1'), (2, 'tag2', 'description2')],
            }
        }

class DBMessage2(BaseModel):
    total_count: int = pydantic.Field(..., description="Question text")
    total_pages: int = pydantic.Field(..., description="Response text")
    dbresult: list[dict[str,str]] = pydantic.Field(..., description="History text")
    count: int = pydantic.Field(0, description="自构建数据集数量")

    class Config:
        schema_extra = {
            "example": {
                "total_count": 1,
                "total_pages": 2,
                "history": [(1, 'tag1', 'description1'), (2, 'tag2', 'description2')],
            }
        }


def get_kb_path(local_doc_id: str):
    return os.path.join(KB_ROOT_PATH, local_doc_id)


def get_doc_path(local_doc_id: str):
    return os.path.join(get_kb_path(local_doc_id), "content")


def get_vs_path(local_doc_id: str):
    return os.path.join(get_kb_path(local_doc_id), "vector_store")


def get_file_path(local_doc_id: str, doc_name: str):
    return os.path.join(get_doc_path(local_doc_id), doc_name)


def validate_kb_name(knowledge_base_id: str) -> bool:
    # 检查是否包含预期外的字符或路径攻击关键字
    if "../" in knowledge_base_id:
        return False
    return True


async def upload_file(
        file: UploadFile = File(description="A single binary file"),
        knowledge_base_id: str = Form(..., description="Knowledge Base Name", example="kb1"),
):
    if not validate_kb_name(knowledge_base_id):
        return BaseResponse(code=403, msg="Don't attack me", data=[])

    saved_path = get_doc_path(knowledge_base_id)
    if not os.path.exists(saved_path):
        os.makedirs(saved_path)

    file_content = await file.read()  # 读取上传文件的内容

    file_path = os.path.join(saved_path, file.filename)
    if os.path.exists(file_path) and os.path.getsize(file_path) == len(file_content):
        file_status = f"文件 {file.filename} 已存在。"
        return BaseResponse(code=200, msg=file_status)

    with open(file_path, "wb") as f:
        f.write(file_content)

    vs_path = get_vs_path(knowledge_base_id)
    vs_path, loaded_files = local_doc_qa.init_knowledge_vector_store([file_path], vs_path)
    if len(loaded_files) > 0:
        file_status = f"文件 {file.filename} 已上传至新的知识库，并已加载知识库，请开始提问。"
        return BaseResponse(code=200, msg=file_status)
    else:
        file_status = "文件上传失败，请重新上传"
        return BaseResponse(code=500, msg=file_status)


async def upload_files(
        files: Annotated[
            List[UploadFile], File(description="Multiple files as UploadFile")
        ],
        knowledge_base_id: str = Form(..., description="Knowledge Base Name", example="kb1"),
):
    if not validate_kb_name(knowledge_base_id):
        return BaseResponse(code=403, msg="Don't attack me", data=[])

    saved_path = get_doc_path(knowledge_base_id)
    if not os.path.exists(saved_path):
        os.makedirs(saved_path)
    filelist = []
    for file in files:
        file_content = ''
        file_path = os.path.join(saved_path, file.filename)
        file_content = await file.read()
        if os.path.exists(file_path) and os.path.getsize(file_path) == len(file_content):
            continue
        with open(file_path, "wb") as f:
            f.write(file_content)
        filelist.append(file_path)
    if filelist:
        vs_path = get_vs_path(knowledge_base_id)
        vs_path, loaded_files = local_doc_qa.init_knowledge_vector_store(filelist, vs_path)
        if len(loaded_files):
            file_status = f"documents {', '.join([os.path.split(i)[-1] for i in loaded_files])} upload success"
            return BaseResponse(code=200, msg=file_status)
    file_status = f"documents {', '.join([os.path.split(i)[-1] for i in loaded_files])} upload fail"
    return BaseResponse(code=500, msg=file_status)


async def list_kbs():
    # Get List of Knowledge Base
    if not os.path.exists(KB_ROOT_PATH):
        all_doc_ids = []
    else:
        all_doc_ids = [
            folder
            for folder in os.listdir(KB_ROOT_PATH)
            if os.path.isdir(os.path.join(KB_ROOT_PATH, folder))
               and os.path.exists(os.path.join(KB_ROOT_PATH, folder, "vector_store", "index.faiss"))
        ]

    return ListDocsResponse(data=all_doc_ids)


import chardet

def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        raw_data = file.read()
        result = chardet.detect(raw_data)
        return result['encoding']

#知识库文件分析
def count_words_in_pdf(file_path):
    try:
        pdf_document = fitz.open(file_path)
        total_words = 0
        for page in pdf_document:
            total_words += len(page.get_text("text").split())
        return total_words
    except:
        return 0

def count_words_in_excel(file_path):
    try:
        workbook = load_workbook(filename=file_path)
        total_words = 0
        for sheet in workbook:
            for row in sheet.iter_rows(values_only=True):
                for cell_value in row:
                    if cell_value and isinstance(cell_value, str):
                        total_words += len(cell_value.split())
        return total_words
    except:
        return 0

def count_words_in_ppt(file_path):
    try:
        presentation = Presentation(file_path)
        total_words = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_words += len(shape.text.split())
        return total_words
    except:
        return 0

def count_words_in_txt(file_path):
    try:
        encoding = detect_encoding(file_path)
        with open(file_path, 'r', encoding=encoding) as file:
            text = file.read()
            total_words = len(text.split())
            return total_words
    except:
        return 0

def count_words_in_word(file_path):
    try:
        doc = docx.Document(file_path)
        total_words = 0
        for paragraph in doc.paragraphs:
            total_words += len(paragraph.text.split())
        return total_words
    except:
        return 0

from apscheduler.schedulers.background import BackgroundScheduler
redisdb_doc = redis.Redis(host="dm.redis.hubpd-internal.com", db = 1, port = 6379, password = 'gX92t@Yv3Xo')

def job():
    all_doc_ids = [
        folder
        for folder in os.listdir(KB_ROOT_PATH)
        if os.path.isdir(os.path.join(KB_ROOT_PATH, folder))
           and os.path.exists(os.path.join(KB_ROOT_PATH, folder, "vector_store", "index.faiss"))
    ]

    for knowledge_base_id in all_doc_ids:
        if not validate_kb_name(knowledge_base_id):
            continue

        knowledge_base_id = urllib.parse.unquote(knowledge_base_id)
        kb_path = get_kb_path(knowledge_base_id)
        if not os.path.exists(kb_path):
            continue

        local_doc_folder = get_doc_path(knowledge_base_id)

        pdf_count = 0
        excel_count = 0
        ppt_count = 0
        txt_count = 0
        docx_count = 0
        word_count = 0

        try:
            for root, dirs, files in os.walk(local_doc_folder):
                for file in files:
                    file_path = os.path.join(root, file)
                    print(file_path)
                    file_extension = file.split('.')[-1].lower()

                    if file_extension == 'pdf':
                        pdf_count += 1
                        word_count += count_words_in_pdf(file_path)
                    elif file_extension in ['xlsx', 'xls']:
                        excel_count += 1
                        word_count += count_words_in_excel(file_path)
                    elif file_extension == 'pptx':
                        ppt_count += 1
                        word_count += count_words_in_ppt(file_path)
                    elif file_extension == 'txt':
                        txt_count += 1
                        word_count += count_words_in_txt(file_path)
                    elif file_extension == 'docx':
                        docx_count += 1
                        word_count += count_words_in_word(file_path)

            konwleage_detail={
                'pdf_count': pdf_count,
                'excel_count': excel_count,
                'ppt_count': ppt_count,
                'txt_count': txt_count,
                'word_count': word_count,
                'docx_count': docx_count,
                'all_file_count': pdf_count + excel_count + ppt_count + txt_count + docx_count
            }

            konwleage_detail_string = json.dumps(konwleage_detail)
            redisdb_doc.set(knowledge_base_id,konwleage_detail_string)

        except Exception as e:
            logging.warning(e)

job()
scheduler = BackgroundScheduler()
scheduler.add_job(job, trigger=IntervalTrigger(hours=2), id='scheduler_doc_id')
scheduler.start()



async def count_files_in_folder(
        knowledge_base_id: str = Query(..., description="Knowledge Base Name", example="kb1")
):
    try:
        konwleage_detail_string = redisdb_doc.get(knowledge_base_id)
        konwleage_json = json.loads(konwleage_detail_string)
        return UserDetailResponse(user=konwleage_json)
    except Exception as e:
        return BaseResponse(code=400,msg=e)



async def get_file_modification_time(file_path):
    try:
        modification_time = os.path.getmtime(file_path)
        modification_datetime = datetime.fromtimestamp(modification_time)
        changetime = modification_datetime.strftime("%Y-%m-%d %H:%M:%S")
        print(changetime)
        return changetime
    except OSError as e:
        print(f"Error: {e}")
        return None


async def list_docs(
        knowledge_base_id: str = Query(..., description="Knowledge Base Name", example="kb1")
):
    if not validate_kb_name(knowledge_base_id):
        return ListDocsResponse(code=403, msg="Don't attack me", data=[])

    knowledge_base_id = urllib.parse.unquote(knowledge_base_id)
    kb_path = get_kb_path(knowledge_base_id)
    local_doc_folder = get_doc_path(knowledge_base_id)
    if not os.path.exists(kb_path):
        return ListDocsResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found", data=[])

    num_files = 0
    if not os.path.exists(local_doc_folder):
        all_doc_names = []
    else:
        files = os.listdir(local_doc_folder)
        num_files = len(files)

        all_doc_names = [
            {
                "name": str(doc),
                "create_time": str(await get_file_modification_time(os.path.join(local_doc_folder, doc))),
                "url": f"http://{host}:8083/local_doc_qa/get_file" + str(os.path.join(local_doc_folder, doc)),
            }
            for doc in os.listdir(local_doc_folder)
            if os.path.isfile(os.path.join(local_doc_folder, doc))
        ]
    return ListDocsCreateResponse(count=num_files, data=all_doc_names)


async def delete_kb(
        knowledge_base_id: str = Query(...,
                                       description="Knowledge Base Name",
                                       example="kb1"),
):
    if not validate_kb_name(knowledge_base_id):
        return BaseResponse(code=403, msg="Don't attack me")

    # TODO: 确认是否支持批量删除知识库
    knowledge_base_id = urllib.parse.unquote(knowledge_base_id)
    kb_path = get_kb_path(knowledge_base_id)
    if not os.path.exists(kb_path):
        return BaseResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found")
    shutil.rmtree(kb_path)
    return BaseResponse(code=200, msg=f"Knowledge Base {knowledge_base_id} delete success")


async def delete_doc(
        knowledge_base_id: str = Query(...,
                                       description="Knowledge Base Name",
                                       example="kb1"),
        doc_name: str = Query(
            ..., description="doc name", example="doc_name_1.pdf"
        ),
):
    if not validate_kb_name(knowledge_base_id):
        return BaseResponse(code=403, msg="Don't attack me")

    knowledge_base_id = urllib.parse.unquote(knowledge_base_id)
    if not os.path.exists(get_kb_path(knowledge_base_id)):
        return BaseResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found")
    doc_path = get_file_path(knowledge_base_id, doc_name)
    if os.path.exists(doc_path):
        os.remove(doc_path)
        remain_docs = await list_docs(knowledge_base_id)
        if len(remain_docs.data) == 0:
            shutil.rmtree(get_kb_path(knowledge_base_id), ignore_errors=True)
            return BaseResponse(code=200, msg=f"document {doc_name} delete success")
        else:
            status = local_doc_qa.delete_file_from_vector_store(doc_path, get_vs_path(knowledge_base_id))
            print(status)
            if "success" in status:
                return BaseResponse(code=200, msg=f"document {doc_name} delete success")
            else:
                return BaseResponse(code=500, msg=f"document {doc_name} delete fail")
    else:
        return BaseResponse(code=404, msg=f"document {doc_name} not found")


async def update_doc(
        knowledge_base_id: str = Query(...,
                                       description="知识库名",
                                       example="kb1"),
        old_doc: str = Query(
            ..., description="待删除文件名，已存储在知识库中", example="doc_name_1.pdf"
        ),
        new_doc: UploadFile = File(description="待上传文件"),
):
    if not validate_kb_name(knowledge_base_id):
        return BaseResponse(code=403, msg="Don't attack me")

    knowledge_base_id = urllib.parse.unquote(knowledge_base_id)
    if not os.path.exists(get_kb_path(knowledge_base_id)):
        return BaseResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found")
    doc_path = get_file_path(knowledge_base_id, old_doc)
    if not os.path.exists(doc_path):
        return BaseResponse(code=404, msg=f"document {old_doc} not found")
    else:
        os.remove(doc_path)
        delete_status = local_doc_qa.delete_file_from_vector_store(doc_path, get_vs_path(knowledge_base_id))
        if "fail" in delete_status:
            return BaseResponse(code=500, msg=f"document {old_doc} delete failed")
        else:
            saved_path = get_doc_path(knowledge_base_id)
            if not os.path.exists(saved_path):
                os.makedirs(saved_path)

            file_content = await new_doc.read()  # 读取上传文件的内容

            file_path = os.path.join(saved_path, new_doc.filename)
            if os.path.exists(file_path) and os.path.getsize(file_path) == len(file_content):
                file_status = f"document {new_doc.filename} already exists"
                return BaseResponse(code=200, msg=file_status)

            with open(file_path, "wb") as f:
                f.write(file_content)

            vs_path = get_vs_path(knowledge_base_id)
            vs_path, loaded_files = local_doc_qa.init_knowledge_vector_store([file_path], vs_path)
            if len(loaded_files) > 0:
                file_status = f"document {old_doc} delete and document {new_doc.filename} upload success"
                return BaseResponse(code=200, msg=file_status)
            else:
                file_status = f"document {old_doc} success but document {new_doc.filename} upload fail"
                return BaseResponse(code=500, msg=file_status)


async def local_doc_chat(
        knowledge_base_id: str = Body(..., description="Knowledge Base Name", example="kb1"),
        question: str = Body(..., description="Question", example="工伤保险是什么？"),
        streaming: bool = Body(False, description="是否开启流式输出，默认false，有些模型可能不支持。"),
        history: List[List[Optional[str]]] = Body(
            [],
            description="History of previous questions and answers",
            example=[
                [
                    "工伤保险是什么？",
                    "工伤保险是指用人单位按照国家规定，为本单位的职工和用人单位的其他人员，缴纳工伤保险费，由保险机构按照国家规定的标准，给予工伤保险待遇的社会保险制度。",
                ]
            ],
        ),
):
    vs_path = get_vs_path(knowledge_base_id)
    if not os.path.exists(vs_path):
        # return BaseResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found")
        return ChatMessage(
            question=question,
            response=f"Knowledge base {knowledge_base_id} not found",
            history=history,
            source_documents=[],
        )
    else:
        if (streaming):
            def generate_answer():
                last_print_len = 0
                for resp, next_history in local_doc_qa.get_knowledge_based_answer(
                        query=question, vs_path=vs_path, chat_history=history, streaming=True
                ):
                    yield resp["result"][last_print_len:]
                    last_print_len = len(resp["result"])

            return StreamingResponse(generate_answer(), media_type="text/event-stream")
        else:
            for resp, history in local_doc_qa.get_knowledge_based_answer(
                    query=question, vs_path=vs_path, chat_history=history, streaming=True
            ):
                print(resp)
                pass

            source_documents = [
                # f"""出处 [{inum + 1}] {os.path.split(doc.metadata['source'])[-1]}：\n\n{doc.page_content}\n\n"""
                # f"""相关度：{doc.metadata['score']}\n\n"""

                {
                    "source": str(os.path.basename(doc.metadata['source'])),
                    "title": str(os.path.basename(doc.metadata['source']).split('.')[0]),
                    "url": f"http://{host}:8083/local_doc_qa/get_file" + str(doc.metadata['source']),
                    "page_content": str(doc.page_content),
                    "score": str(doc.metadata['score'])
                }

                for inum, doc in enumerate(resp["source_documents"])
            ]

            return ChatMessage(
                question=question,
                response=resp["result"],
                history=history,
                source_documents=source_documents,
            )


async def chat_v2(
        question: str = Body(..., description="Question", example="工伤保险是什么？"),
        streaming: bool = Body(False, description="是否开启流式输出，默认false，有些模型可能不支持。"),
        history: List[List[Optional[str]]] = Body(
            [],
            description="History of previous questions and answers",
            example=[
                [
                    "工伤保险是什么？",
                    "工伤保险是指用人单位按照国家规定，为本单位的职工和用人单位的其他人员，缴纳工伤保险费，由保险机构按照国家规定的标准，给予工伤保险待遇的社会保险制度。",
                ]
            ],
        ),
):
    vs_path = get_vs_path("leader")
    if not os.path.exists(vs_path):
        # return BaseResponse(code=404, msg=f"Knowledge base {knowledge_base_id} not found")
        return ChatMessage(
            question=question,
            response=f"Knowledge base  not found",
            history=history,
            source_documents=[],
        )
    else:
        if (streaming):
            def generate_answer():
                last_print_len = 0
                for resp, next_history in local_doc_qa.get_knowledge_based_answer(
                        query=question, vs_path=vs_path, chat_history=history, streaming=True
                ):
                    yield resp["result"][last_print_len:]
                    last_print_len = len(resp["result"])

            return StreamingResponse(generate_answer())
        else:
            for resp, history in local_doc_qa.get_knowledge_based_answer(
                    query=question, vs_path=vs_path, chat_history=history, streaming=True
            ):
                # print(resp)
                pass

            source_documents = []

            return ChatMessage(
                question=question,
                response=resp["result"],
                history=history,
                source_documents=source_documents,
            )


async def bing_search_chat(
        question: str = Body(..., description="Question", example="工伤保险是什么？"),
        history: Optional[List[List[Optional[str]]]] = Body(
            [],
            description="History of previous questions and answers",
            example=[
                [
                    "工伤保险是什么？",
                    "工伤保险是指用人单位按照国家规定，为本单位的职工和用人单位的其他人员，缴纳工伤保险费，由保险机构按照国家规定的标准，给予工伤保险待遇的社会保险制度。",
                ]
            ],
        ),
):
    for resp, history in local_doc_qa.get_search_result_based_answer(
            query=question, chat_history=history, streaming=True
    ):
        print(resp)
        pass

    source_documents = [
        # f"""出处 [{inum + 1}] [{doc.metadata["source"]}]({doc.metadata["source"]}) \n\n{doc.metadata["filename"]} \n\n{doc.page_content}\n\n"""

        {
            "url": str(doc.metadata['source']),
            "title": str(doc.metadata['filename']),
            "page_content": str(doc.page_content)
        }

        for inum, doc in enumerate(resp["source_documents"])
    ]

    return ChatMessage(
        question=question,
        response=resp["result"],
        history=history,
        source_documents=source_documents,
    )


async def chat(
        question: str = Body(..., description="Question", example="工伤保险是什么？"),
        streaming: bool = Body(False, description="是否开启流式输出，默认false，有些模型可能不支持。"),
        history: List[List[Optional[str]]] = Body(
            [],
            description="History of previous questions and answers",
            example=[
                [
                    "工伤保险是什么？",
                    "工伤保险是指用人单位按照国家规定，为本单位的职工和用人单位的其他人员，缴纳工伤保险费，由保险机构按照国家规定的标准，给予工伤保险待遇的社会保险制度。",
                ]
            ],
        ),
):
    if (streaming):
        def generate_answer():
            last_print_len = 0
            answer_result_stream_result = local_doc_qa.llm_model_chain(
                {"prompt": question, "history": history, "streaming": True})
            for answer_result in answer_result_stream_result['answer_result_stream']:
                yield answer_result.llm_output["answer"][last_print_len:]
                last_print_len = len(answer_result.llm_output["answer"])

        return StreamingResponse(generate_answer(), media_type="text/event-stream")
    else:
        answer_result_stream_result = local_doc_qa.llm_model_chain(
            {"prompt": question, "history": history, "streaming": True})
        for answer_result in answer_result_stream_result['answer_result_stream']:
            resp = answer_result.llm_output["answer"]
            history = answer_result.history
            pass

        return ChatMessage(
            question=question,
            response=resp,
            history=history,
            source_documents=[],
        )
    answer_result_stream_result = local_doc_qa.llm_model_chain(
        {"prompt": question, "history": history, "streaming": True})

    for answer_result in answer_result_stream_result['answer_result_stream']:
        resp = answer_result.llm_output["answer"]
        history = answer_result.history
        pass
    return ChatMessage(
        question=question,
        response=resp,
        history=history,
        source_documents=[],
    )


async def stream_chat(websocket: WebSocket):
    await websocket.accept()
    turn = 1
    while True:
        input_json = await websocket.receive_json()
        question, history, knowledge_base_id = input_json["question"], input_json["history"], input_json[
            "knowledge_base_id"]
        vs_path = get_vs_path(knowledge_base_id)

        if not os.path.exists(vs_path):
            await websocket.send_json({"error": f"Knowledge base {knowledge_base_id} not found"})
            await websocket.close()
            return

        await websocket.send_json({"question": question, "turn": turn, "flag": "start"})

        last_print_len = 0
        for resp, history in local_doc_qa.get_knowledge_based_answer(
                query=question, vs_path=vs_path, chat_history=history, streaming=True
        ):
            await asyncio.sleep(0)
            await websocket.send_text(resp["result"][last_print_len:])
            last_print_len = len(resp["result"])

        source_documents = [
            f"""出处 [{inum + 1}] {os.path.split(doc.metadata['source'])[-1]}：\n\n{doc.page_content}\n\n"""
            f"""相关度：{doc.metadata['score']}\n\n"""
            for inum, doc in enumerate(resp["source_documents"])
        ]

        await websocket.send_text(
            json.dumps(
                {
                    "question": question,
                    "turn": turn,
                    "flag": "end",
                    "sources_documents": source_documents,
                },
                ensure_ascii=False,
            )
        )
        turn += 1


async def stream_chat_bing(websocket: WebSocket):
    """
    基于bing搜索的流式问答
    """
    await websocket.accept()
    turn = 1
    while True:
        input_json = await websocket.receive_json()
        question, history = input_json["question"], input_json["history"]

        await websocket.send_json({"question": question, "turn": turn, "flag": "start"})

        last_print_len = 0
        for resp, history in local_doc_qa.get_search_result_based_answer(question, chat_history=history,
                                                                         streaming=True):
            await websocket.send_text(resp["result"][last_print_len:])
            last_print_len = len(resp["result"])

        source_documents = [
            f"""出处 [{inum + 1}] {os.path.split(doc.metadata['source'])[-1]}：\n\n{doc.page_content}\n\n"""
            f"""相关度：{doc.metadata['score']}\n\n"""
            for inum, doc in enumerate(resp["source_documents"])
        ]

        await websocket.send_text(
            json.dumps(
                {
                    "question": question,
                    "turn": turn,
                    "flag": "end",
                    "sources_documents": source_documents,
                },
                ensure_ascii=False,
            )
        )
        turn += 1


async def document():
    return RedirectResponse(url="/docs")


async def searchdb(user_name):
    select_sql = "select user_id,user_name,password,knowledge_base_id,perms,phone_num,create_time from t_knowledge_user_info where user_name = '%s'" % (
        user_name)

    # conn = pymysql.connect(host='10.19.1.1'  # 连接名称，默认127.0.0.1
    #                        , user='theshy'  # 用户名
    #                        , passwd='123456Mysql!'  # 密码
    #                        , port=3306  # 端口，默认为3306
    #                        , db='langchain'  # 数据库名称
    #                        , charset='utf8'
    #                        , autocommit=True,  # 如果插入数据，自动提交给数据库
    #                        )

    conn = pymysql.connect(host='pdmi-db-1.rwlb.rds.aliyuncs.com'  # 连接名称，默认127.0.0.1
                           , user='llm'  # 用户名
                           , passwd='EbcLvd2=YLL*'  # 密码
                           , port=3306  # 端口，默认为3306
                           , db='llm'  # 数据库名称
                           , charset='utf8'
                           , autocommit=True,  # 如果插入数据，自动提交给数据库
                           )

    cur = conn.cursor()
    result = []

    try:
        cur.execute(select_sql)
        conn.commit()  # 提交到数据库执行
        result = cur.fetchall()
    except Exception as w:
        conn.rollback()
    finally:
        conn.close()

    return result


# 用户登录
async def login(
        user_name: str = Query(...,
                               description="用户姓名",
                               example="kb1"),
        password: str = Query(
            ..., description="用户密码", example="123456"
        ),
):
    import json
    userexist = redisdb.exists(user_name)
    password_db = ""
    results = []
    result = []

    if userexist:
        usesr_detail = redisdb.get(user_name)
        user_json = json.loads(usesr_detail)
        password_db = user_json["password"]
    else:
        results = await searchdb(user_name)
        if len(results) == 0:
            return BaseResponse(code=404, msg=f"user {user_name} haved not sign in")
        else:
            result = results[0]
            password_db = result[2]

    if password != password_db:
        return BaseResponse(code=404, msg=f"password error!")
    else:
        if not userexist:
            user_json = {"user_id": result[0],
                         "user_name": result[1],
                         "password": result[2],
                         "knowledge_base_id": result[3],
                         "perms": result[4],
                         "phone_num": result[5],
                         "create_time": result[6]}

            detail_string = json.dumps(user_json)
            redisdb.set(user_name, detail_string)

        from datetime import datetime, timedelta
        payload = {
            'exp': datetime.now() + timedelta(hours=1),  # 令牌过期时间
            'username': user_name  # 想要传递的信息,如用户名ID
        }
        key = 'LANG_CHAIN'

        encoded_jwt = jwt.encode(payload, key, algorithm='HS256')
        return BaseResponse(code=200, msg=encoded_jwt)


# 获取用户信息列表
async def user_detail(
        token: str = Query(...,
                           description="用户登录token",
                           example="D9TC4UbxuyTKow1GH6OeXFjiPEkH54w"),

):
    try:
        jwt_decode = jwt.decode(token, 'LANG_CHAIN', algorithms=['HS256'])
    except Exception as e:
        return BaseResponse(code=404, msg=f"the token {token} is not normalize")

    # user_decode = json.loads(jwt_decode)
    username = jwt_decode["username"]

    userexist = redisdb.exists(username)
    if userexist:
        user_detail = redisdb.get(username)
        user_json = json.loads(user_detail)
        return UserDetailResponse(user=user_json)
    else:
        results = await searchdb(username)

        if len(results) == 0:
            return BaseResponse(code=404, msg=f"The user corresponding to the token {token} does not exist")
        else:
            result = results[0]

            user_json = {"user_id": result[0],
                         "user_name": result[1],
                         "password": result[2],
                         "knowledge_base_id": result[3],
                         "perms": result[4],
                         "phone_num": result[5],
                         "create_time": result[6]}

            detail_string = json.dumps(user_json)
            redisdb.set(username, detail_string)

            return UserDetailResponse(user=user_json)


async def logout(
        token: str = Query(...,
                           description="用户登录token",
                           example="D9TC4UbxuyTKow1GH6OeXFjiPEkH54w"),

):
    jwt_decode = jwt.decode(token, 'LANG_CHAIN', algorithms=['HS256'])
    username = jwt_decode["username"]

    userexist = redisdb.exists(username)
    if userexist:
        user_detail = redisdb.get(username)
        return BaseResponse(code=200, msg=user_detail)
    else:
        return BaseResponse(code=404, msg=f"The user corresponding to the token {token} does not exist")


async def get_file(file_path):
    if os.path.exists(file_path):
        file_name = os.path.basename(file_path)
        utf8_encoded_filename = file_name.encode("utf-8").decode("latin1")
        print(utf8_encoded_filename)
        # return FileResponse(file_path, headers={"Content-Disposition": f"attachment; filename={utf8_encoded_filename}"})

        return FileResponse(file_path, filename=utf8_encoded_filename)

    else:
        return {"error": "File not found"}




#复制数据集接口
def copy_dataset_record(dataset_id, new_name):
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 查询原始数据
            select_sql = "SELECT * FROM t_dataset WHERE dataset_id = %s"
            cursor.execute(select_sql, (dataset_id,))
            original_data = cursor.fetchone()

            if original_data:
                # 提取需要复制的字段
                source_dataset_id = original_data[0]
                source_dataset_name = original_data[5]
                insert_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

                # 插入新数据
                insert_sql = """
                    INSERT INTO t_dataset (dataset_id, sorting, source_platform, platform_url, platform_type, 
                                           name, name_en, summary, description, download_url, type, task_type, tags, 
                                           domain, purpose, region, format, size, language, license, publisher, 
                                           preview_url, outbound_link, publicly, remark, data_area, download_count, 
                                           favorites_count, source_datasetid, source_datasetname, governance, 
                                           insert_time, update_time)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, 
                            %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                """
                cursor.execute(insert_sql, (
                    0, original_data[1], original_data[2], original_data[3], original_data[4],
                    new_name, original_data[6], original_data[7], original_data[8], original_data[9],
                    original_data[10], original_data[11], original_data[12], original_data[13], original_data[14],
                    original_data[15], original_data[16], original_data[17], original_data[18], original_data[19],
                    original_data[20], original_data[21], original_data[22], original_data[23], original_data[24],
                    original_data[25], original_data[26], original_data[27], original_data[28],
                    source_dataset_id, source_dataset_name, original_data[30], insert_time, None
                ))
                conn.commit()

                return "新数据已成功写入"
            else:
                return "未找到符合条件的原始数据"

    finally:
        # 释放连接
        conn.close()


import json



async def get_dataset_data2(

    tags: str = Query(...,
                           description="查询字段",
                           example=""),
):
    json_tagsq = json.loads(tags)
    # 从连接池获取连接
    conn = pool.connection()
    page = 1
    page_size = 10
    json_tags = {}
    for key, value in json_tagsq.items():
        if key == "page" :
            page = json_tagsq["page"]
        elif key == "page_size" :
            page_size = json_tagsq["page_size"]
        elif value != "":
            json_tags[key] = value


    try:
        with conn.cursor() as cursor:
            # 构建查询语句
            conditions = []


            if "dataset_id" in json_tags:
                darasetid = json_tags["dataset_id"]
                sql = """
                                SELECT *
                                FROM t_dataset
                                WHERE dataset_id = %s
                            """% (darasetid)

            else:
                # 解析tags中的条件
                for key, value in json_tags.items():
                    if key == "keyword":
                        conditions.append(f"(name LIKE '%{value}%' or summary LIKE '%{value}%')")
                    else:
                        conditions.append(f"{key} LIKE '%{value}%'")



                sql = """
                    SELECT *
                    FROM t_dataset
                    WHERE {}
                """.format(" AND ".join(conditions))



                sql= sql+ f" order by sorting asc, insert_time desc"

            # 查询总数
            cursor.execute(f"SELECT COUNT(*) FROM ({sql}) AS count_query")
            total_count = cursor.fetchone()[0]

            # 计算分页
            total_pages = (total_count + page_size - 1) // page_size

            # 执行分页查询
            start_index = (page - 1) * page_size
            sql += " LIMIT %s, %s"% (start_index, page_size)

            cursor.execute(sql,)
            result = cursor.fetchall()

            field_names = [i[0] for i in cursor.description]
            lists = []
            for i in result:
                list = {}
                for index,r in enumerate(i):
                    if r is None:
                        list[field_names[index]] = "null"
                    else:
                        list[field_names[index]] = str(r)

                lists.append(list)

            if "data_area" in json_tags:
                if json_tags["data_area"] == "自建":
                    sql = """
                                       SELECT count(1)
                                       FROM t_dataset
                                       WHERE data_area = "自建"
                                   """

                    # 查询总数
                    cursor.execute(sql)
                    zijian_count = cursor.fetchone()[0]
                    return DBMessage(total_count=total_count, total_pages=total_pages,
                                     dbresult=lists, count=zijian_count)

            return DBMessage( total_count=total_count, total_pages=total_pages,
                            dbresult=lists)

    except Exception as e:
        return BaseResponse(code=404, msg=str(e))

    finally:
        # 释放连接
        conn.close()





async def get_dataset_data(
    page: int = Form(1),
    page_size: int = Form(10)
):

    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 构建查询语句
            conditions = []
            param = []


            sql = """
                SELECT *
                FROM t_dataset
                WHERE data_area = "公开"
            """


            # 查询总数
            cursor.execute(f"SELECT COUNT(*) FROM ({sql}) AS count_query")
            total_count = cursor.fetchone()[0]

            # 计算分页
            total_pages = (total_count + page_size - 1) // page_size

            # 执行分页查询
            start_index = (page - 1) * page_size
            sql += " LIMIT %s, %s"
            param += (start_index, page_size)
            cursor.execute(sql, param)
            result = cursor.fetchall()

            field_names = [i[0] for i in cursor.description]
            lists = []
            for i in result:
                list = {}
                for index,r in enumerate(i):
                    if r is None:
                        print("dd")
                        list[field_names[index]] = "null"
                    else:
                        list[field_names[index]] = str(r)

                print(list)
                lists.append(list)

            return DBMessage( total_count=total_count, total_pages=total_pages,
                            dbresult=lists)

    except Exception as e:
        return BaseResponse(code=404, msg=str(e))

    finally:
        # 释放连接
        conn.close()




async def get_preview_urls(
        page_size: int = Body(...,
                              description="查询字段",
                              example=""),
        page: int = Body(...,
                         description="查询字段",
                         example=""),

        dataset_id: int = Body(...,
                               description="查询字段",
                               example=""),
        keyword: Optional[str] = None
):
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 执行查询

            if keyword is None or keyword == "":
                sql = "SELECT name, Summary, description, size, preview_url,format,insert_time FROM t_dataset_file WHERE dataset_id=%s"%(dataset_id)
            else:
                sql = "SELECT name, Summary, description, size, preview_url,format,insert_time FROM t_dataset_file WHERE dataset_id=%s and (name like '%s' or Summary like '%s' or description like '%s')" % (dataset_id,f'%{keyword}%', f'%{keyword}%', f'%{keyword}%')

            cursor.execute(f"SELECT COUNT(*) FROM ({sql}) AS count_query")
            total_count = cursor.fetchone()[0]

            # 计算分页
            total_pages = (total_count + page_size - 1) // page_size

            # 执行分页查询
            start_index = (page - 1) * page_size
            sql += f" LIMIT %s, %s" % (start_index, page_size)
            cursor.execute(sql, )
            result = cursor.fetchall()

            field_names = [i[0] for i in cursor.description]
            lists = []
            for i in result:
                list = {}
                for index, r in enumerate(i):
                    if r is None:
                        list[field_names[index]] = "null"
                    else:
                        list[field_names[index]] = str(r)

                print(list)
                lists.append(list)

            return DBMessage(total_count=total_count,
                             total_pages=total_pages, dbresult=lists)
    except Exception as e:
            return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()



async def get_dataset_files(

        dataset_id: int = Body(...,
                               description="查询字段",
                               example=""),
        dataset_name: str = Body(...,
                               description="查询字段",
                               example=""),
):
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 执行查询

            sql = "SELECT preview_url FROM t_dataset_file WHERE dataset_id=%s" %(dataset_id)
            cursor.execute(sql, )
            result = cursor.fetchall()
            print(result)


            lists = []
            for index, item in enumerate(result):
                lists.append(item[0])

            return ListDocsResponse(data=lists)
    except Exception as e:
            return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()




async def insert_data_value(
        prompt: str = Body(...,
                          description="查询字段",
                          example=""),

        context: Optional[str] = None,

        response: str = Body(...,
                          description="查询字段",
                          example=""),

        dataset_id: str = Body(...,
                          description="查询字段",
                          example=""),


):
    # 从连接池获取连接
    conn = pool.connection()
    if(context==None):
        context = "null"

    try:
        with conn.cursor() as cursor:
            # 执行插入操作
            sql = "INSERT INTO t_data_values (prompt, context, response, dataset_id, publicly, remark, insert_time) VALUES (%s, %s, %s, %s, 'N', '', %s)"
            insert_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            cursor.execute(sql, (prompt, context, response, dataset_id, insert_time))

        # 提交事务
        conn.commit()
        return BaseResponse(code=200,msg="create success")
    except Exception as e:
        return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()


async def insert_data_novalues(
        prompt: str = Body(...,
                          description="查询字段",
                          example=""),

        context: Optional[str] = None,

        response: str = Body(...,
                          description="查询字段",
                          example=""),

        dataset_id: str = Body(...,
                          description="查询字段",
                          example=""),

):
    # 从连接池获取连接
    conn = pool.connection()
    if(context==None):
        context = "null"

    try:
        with conn.cursor() as cursor:
            # 执行插入操作
            sql = "INSERT INTO t_data_novalues (prompt, context, response, dataset_id, publicly, remark, insert_time) VALUES (%s, %s, %s, %s, 'N', '', %s)"
            insert_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            cursor.execute(sql, (prompt, context, response, dataset_id, insert_time))

        # 提交事务
        conn.commit()
        return BaseResponse(code=200,msg="create success")
    except Exception as e:
        return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()


async def insert_dataset(
    files: List[str] = Body(...,description="查询字段",example=""),
    name_en: str = Body("default dataset name", description="查询字段", example=""),
    name: str = Body(...,description="查询字段",example=""),
    publisher: str = Body(...,description="查询字段",example=""),
    license: str = Body(...,description="查询字段",example=""),
    tags: str = Body(...,description="查询字段",example=""),
    summary: str = Body(...,description="查询字段",example=""),
    domain: str = Body(None, description="查询字段", example=""),
    purpose: str = Body(None, description="查询字段", example=""),
):
    # 在函数体内使用参数
        # 从连接池获取连接
        conn = pool.connection()

        if(name_en == "" or name_en is None):
            name_en = "default dataset name"

        if ((tags == "" or tags is None) and (domain == "" or domain is None) and (purpose == "" or purpose is None)):
            return BaseResponse(code=400, msg="tags domain purpose 不能均为空")

        try:
            with conn.cursor() as cursor:
                base_url = r'/data/llms/langchain-ChatGLM/dataset/selfbuild/'
                #base_url = r'/Applications/code/pdmi/daimapdmi/langchain-ChatGLM/test/'
                upload_dir = base_url + name

                if not os.path.exists(upload_dir):
                    os.makedirs(upload_dir)
                else:
                    return BaseResponse(code=404, msg="数据集名称已存在")

                preview_url = ""
                format = ""
                insert_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                # 执行插入操作
                sql = """
                    INSERT INTO t_dataset (
                        sorting,source_platform, platform_url, platform_type, name, name_en, summary, description,
                        download_url, type, task_type, tags, domain, purpose, region, format, size, language,
                        license, publisher, preview_url, outbound_link, publicly, remark, data_area,
                        download_count, favorites_count, source_datasetid, source_datasetname, governance,
                        insert_time, update_time
                    ) VALUES(%s, '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', '%s', %s, %s, '%s', '%s', '%s', '%s', '%s')
                    """ %(
                    10000, "自建数据集平台", "12", "1", name, name_en, summary, summary,
                    "1", "type", "task_type", tags, domain, purpose, "1", "1", format, "dataset_size",
                    license, publisher, preview_url, "1", "1",
                    "1", "自建", 0, 0, "1", "1","1",
                    insert_time, insert_time
                )
                cursor.execute(sql)

                # 记录下刚写的数据集id
                last_inserted_id = cursor.lastrowid

                for filepath in files:
                    shutil.move(filepath, upload_dir)
                    filenameo = os.path.basename(filepath)
                    file_path = os.path.join(upload_dir, filenameo)

                    file_size = get_size(file_path)
                    if preview_url == "":
                        preview_url = r'http://dm.hubpd.com:8083/local_doc_qa/get_file' + file_path
                        format = filenameo.split(".")[-1]
                    # 在数据库表 t_dataset_file 中插入一条数据

                    sql = """
                                INSERT INTO t_dataset_file
                                (dataset_id, sorting, name, Summary, description, download_url, type, tags, format, size, preview_url, publicly, remark, download_count, favorites_count, insert_time, update_time)
                                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                            """
                    cursor.execute(sql, (
                        last_inserted_id, 10000, filenameo, "summary", "description", preview_url, type, "tags",
                        format, file_size, preview_url, "publicly",
                        "remark", 0, 0, datetime.now(), None))

                dataset_size = get_size(upload_dir)
                update_sql = "update t_dataset set size = '%s',format = '%s', preview_url = '%s' where dataset_id = '%s'" % (dataset_size,format,preview_url,last_inserted_id)
                cursor.execute(update_sql)
                conn.commit()

            return UIDResponse(last_id=last_inserted_id)
        except Exception as e:
            return BaseResponse(code=404, msg=str(e))

        finally:
            # 释放连接
            conn.close()




async def search_data_values(
        page_size: int = Body(...,
                           description="查询字段",
                           example=""),
        keyword: Optional[str] = None,
        page: int = Body(...,
                             description="查询字段",
                             example=""),
        dataset_id :int= Body(...,
                             description="查询字段",
                             example=""),
):
    #keyword = json_tags["keyword"]
    # 从连接池获取连接
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 构建查询语句
            if keyword is None:
                sql = """
                        SELECT *
                        FROM t_data_values 
                        where dataset_id = %s 
                    """ % (dataset_id)
            else:
                sql = """
                        SELECT *
                        FROM t_data_values 
                        where dataset_id = %s and (prompt LIKE '%s' OR context LIKE '%s' OR response LIKE '%s')
                    """ % (dataset_id, f'%{keyword}%', f'%{keyword}%', f'%{keyword}%')
            # WHERE prompt LIKE %s OR context LIKE %s OR response LIKE %s
            #param = (f"%{keyword}%", f"%{keyword}%", f"%{keyword}%")

            print(sql)

            # 查询总数
            cursor.execute(f"SELECT COUNT(*) FROM ({sql}) AS count_query")
            total_count = cursor.fetchone()[0]

            # 计算分页
            total_pages = (total_count + page_size - 1) // page_size

            # 执行分页查询
            start_index = (page - 1) * page_size
            sql += f" LIMIT %s, %s"%(start_index, page_size)
            cursor.execute(sql,)
            result = cursor.fetchall()

            field_names = [i[0] for i in cursor.description]
            lists = []
            for i in result:
                list = {}
                for index,r in enumerate(i):
                    if r is None:
                        print("dd")
                        list[field_names[index]] = "null"
                    else:
                        list[field_names[index]] = str(r)

                print(list)
                lists.append(list)

            return DBMessage(total_count=total_count,
                             total_pages=total_pages, dbresult=lists)
    except Exception as e:
        return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()


async def search_data_novalues(
        page_size: int = Body(...,
                              description="查询字段",
                              example=""),
        keyword: Optional[str] = None,
        page: int = Body(...,
                         description="查询字段",
                         example=""),
        dataset_id: int = Body(...,
                               description="查询字段",
                               example=""),
):
    # 从连接池获取连接
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 构建查询语句
            if keyword is None:
                sql = """
                        SELECT *
                        FROM t_data_novalues
                        where dataset_id = %s 
                    """ % (dataset_id)
            else:
                sql = """
                        SELECT *
                        FROM t_data_novalues
                        where dataset_id = %s and (prompt LIKE '%s' OR context LIKE '%s' OR response LIKE '%s')
                    """ % (dataset_id, f'%{keyword}%', f'%{keyword}%', f'%{keyword}%')
            # WHERE prompt LIKE %s OR context LIKE %s OR response LIKE %s
            # param = (f"%{keyword}%", f"%{keyword}%", f"%{keyword}%")

            # 查询总数
            cursor.execute(f"SELECT COUNT(*) FROM ({sql}) AS count_query")
            total_count = cursor.fetchone()[0]

            # 计算分页
            total_pages = (total_count + page_size - 1) // page_size

            # 执行分页查询
            start_index = (page - 1) * page_size
            sql += f" LIMIT %s, %s" % (start_index, page_size)
            cursor.execute(sql, )
            result = cursor.fetchall()

            field_names = [i[0] for i in cursor.description]
            lists = []
            for i in result:
                list = {}
                for index, r in enumerate(i):
                    if r is None:
                        print("dd")
                        list[field_names[index]] = "null"
                    else:
                        list[field_names[index]] = str(r)

                print(list)
                lists.append(list)

            return DBMessage(total_count=total_count,
                             total_pages=total_pages, dbresult=lists)
    except Exception as e:
        return BaseResponse(code=404, msg=e)

    finally:
        # 释放连接
        conn.close()




async def upload_dataset_files(
        files: Annotated[
            List[UploadFile], File(description="Multiple files as UploadFile")
        ],
        dataset_id: str = Form(..., description="Knowledge Base Name", example="kb1"),
        dataset_name: str = Form(..., description="Knowledge Base Name", example="kb1"),
):
    #upload_dir = "/data/llms/langchain-ChatGLM/dataset/selfbuild/"+dataset_name
    upload_dir = "/Applications/code/pdmi/daimapdmi/langchain-ChatGLM/test/" + dataset_name
    conn = pool.connection()
    cursor = conn.cursor()
    try:
        search_sql = "select dataset_id,name from t_dataset where dataset_id = %s" %(dataset_id)
        cursor.execute(search_sql)
        result = cursor.fetchone()

        if len(result) == 0:
            return BaseResponse(code=404, msg=f'数据集id{dataset_id}不存在')
        elif result[1] != dataset_name:
            return BaseResponse(code=404, msg=f'数据集id与数据集名称{dataset_name}不符')


        search_sql_count = "select dataset_id,name from t_dataset_file where dataset_id = %s" %(dataset_id)
        cursor.execute(search_sql_count)
        result_count = cursor.fetchone()


        # 检查上传目录是否存在，如果不存在则创建
        if not os.path.exists(upload_dir):
            os.makedirs(upload_dir)

        filelist = []
        preview_url = ""
        format = ""
        for file in files:
            file_path = os.path.join(upload_dir, file.filename)
            file_content = await file.read()

            with open(file_path, "wb") as f:
                f.write(file_content)

            file_size = get_size(file_path)
            filelist.append(file_path)
            preview_url = r'http://dm.hubpd.com:8083/local_doc_qa/get_file'+file_path
            format = file.filename.split(".")[-1]
            # 在数据库表 t_dataset_file 中插入一条数据

            sql = """
                        INSERT INTO t_dataset_file
                        (dataset_id, sorting, name, Summary, description, download_url, type, tags, format, size, preview_url, publicly, remark, download_count, favorites_count, insert_time, update_time)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    """
            cursor.execute(sql, (
            dataset_id, dataset_id, file.filename, "summary", "description", preview_url, type, "tags", format, file_size, preview_url, "publicly",
            "remark", 0, 0, datetime.now(), None))


        dataset_size = get_size(upload_dir)
        if result_count is None:
            update_sql = "update t_dataset set size = '%s',format = '%s', preview_url = '%s' where dataset_id = '%s'" % (dataset_size,format,preview_url,dataset_id)
        else:
            update_sql = "update t_dataset set size = '%s' where dataset_id = '%s'" % (dataset_size, dataset_id)
        cursor.execute(update_sql)
        conn.commit()
        return BaseResponse(code=200, msg="upload success")

    except Exception as e:
        return BaseResponse(code=404, msg="upload fail")
    finally:
        conn.close()



async def upload_dataset_files_temp(
        file: Annotated[
            UploadFile, File(description="Multiple files as UploadFile")
        ],
):

    upload_dir = "/data/llms/langchain-ChatGLM/dataset/selfbuild/temp/"
    #upload_dir = "/Applications/code/pdmi/daimapdmi/langchain-ChatGLM/test/"
    # 检查上传目录是否存在，如果不存在则创建
    if not os.path.exists(upload_dir):
        os.makedirs(upload_dir)

    try:
        file_path = os.path.join(upload_dir, file.filename)
        file_content = await file.read()

        with open(file_path, "wb") as f:
            f.write(file_content)

        return BaseResponse(code=200, msg=file_path)

    except Exception as e:
        return BaseResponse(code=400, msg=e)



async def copyto_data_values(
        dataset_id: int = Body(...,
                               description="查询字段",
                               example=""),
        data_area: str = Body(...,
                               description="查询字段",
                               example=""),
):
    conn = pool.connection()
    print(dataset_id)

    try:
        with conn.cursor() as cursor:
            # 查询原始数据
            select_sql = "SELECT * FROM t_dataset WHERE dataset_id = %s"
            cursor.execute(select_sql, (dataset_id,))
            original_data = cursor.fetchone()
            print(len(original_data))

            if original_data:
                # 提取需要复制的字段
                source_dataset_id = original_data[0]
                sorting = int(original_data[1]) + 1
                source_dataset_name = original_data[5]
                insert_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

                # 插入新数据
                insert_sql = """
                           INSERT INTO t_dataset (dataset_id, sorting, source_platform, platform_url, platform_type, 
                                                  name, name_en, summary, description, download_url, type, task_type, tags, 
                                                  domain, purpose, region, format, size, language, license, publisher, 
                                                  preview_url, outbound_link, publicly, remark, data_area, download_count, 
                                                  favorites_count, source_datasetid, source_datasetname, governance, 
                                                  insert_time, update_time)
                           VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, 
                                   %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                       """

                cursor.execute(insert_sql, (
                    0, sorting, original_data[2], original_data[3], original_data[4],
                    original_data[5], original_data[6], original_data[7], original_data[8], original_data[9],
                    original_data[10], original_data[11], original_data[12], original_data[13], original_data[14],
                    original_data[15], original_data[16], original_data[17], original_data[18], original_data[19],
                    original_data[20], original_data[21], original_data[22], original_data[23], "1",
                     "治理", original_data[26], original_data[27],
                    source_dataset_id, source_dataset_name, original_data[30], insert_time, insert_time
                ))

            update_sql = "update t_dataset set remark = '1' WHERE dataset_id = %s"
            cursor.execute(update_sql, (dataset_id,))

            # 提交事务
            conn.commit()
            return BaseResponse()

    except Exception as e:
        print(e)
        # 发生异常时回滚事务
        conn.rollback()
        return BaseResponse(code=404,msg="fail")
    finally:
        # 关闭连接
        conn.close()

async def edit_data_values(
        prompt: str = Body(...,
                          description="查询字段",
                          example=""),

        context: str = Body(...,
                          description="查询字段",
                          example=""),

        response: str = Body(...,
                          description="查询字段",
                          example=""),

        datavalue_id: str = Body(...,
                          description="查询字段",
                          example=""),
):
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 1. 更新 t_data_values 表
            update_values_sql = """
                UPDATE t_data_values 
                SET prompt = %s, context = %s, response = %s 
                WHERE id = %s
            """
            cursor.execute(update_values_sql, (
                prompt,  # 传入的 prompt 值
                context,  # 传入的 context 值
                response,  # 传入的 response 值
                datavalue_id  # 传入的 dataset_id 值
            ))
            # 提交事务
            conn.commit()
            return BaseResponse(code=200, msg="success")
    except Exception as e:
        # 发生异常时回滚事务
        conn.rollback()
        print(f"Error: {e}")
        return BaseResponse(code=404, msg=e)
    finally:
        # 关闭连接
        conn.close()





async def edit_data_novalues(
        prompt: str = Body(...,
                          description="查询字段",
                          example=""),

        context: str = Body(...,
                          description="查询字段",
                          example=""),

        response: str = Body(...,
                          description="查询字段",
                          example=""),

        datavalue_id: str = Body(...,
                          description="查询字段",
                          example=""),
):
    conn = pool.connection()

    try:
        with conn.cursor() as cursor:
            # 1. 更新 t_data_values 表
            update_values_sql = """
                UPDATE t_data_novalues 
                SET prompt = %s, context = %s, response = %s 
                WHERE id = %s
            """
            cursor.execute(update_values_sql, (
                prompt,  # 传入的 prompt 值
                context,  # 传入的 context 值
                response,  # 传入的 response 值
                datavalue_id  # 传入的 dataset_id 值
            ))
            # 提交事务
            conn.commit()
            return BaseResponse(code=200, msg="success")
    except Exception as e:
        # 发生异常时回滚事务
        conn.rollback()
        print(f"Error: {e}")
        return BaseResponse(code=404, msg=e)
    finally:
        # 关闭连接
        conn.close()


def get_size(path):
    if os.path.isfile(path):
        sizeMb = os.path.getsize(path) / (1024 * 1024)
        return f'{sizeMb:.2f}m'  # 返回文件大小（MB）
    elif os.path.isdir(path):
        total_size = 0
        for dirpath, dirnames, filenames in os.walk(path):
            for f in filenames:
                fp = os.path.join(dirpath, f)
                total_size += os.path.getsize(fp)
        sizeMb = total_size / (1024 * 1024)
        return f'{sizeMb:.2f}m'  # 返回目录大小（MB）
    else:
        return "0m"


async def text2images(
        width : int = Body(...,
                               description="高度",
                               example=""),
        height : int = Body(...,
                               description="宽度",
                               example=""),
        prompt : str = Body(...,
                               description="正向提示语",
                               example=""),
        negative_prompt : Optional[str] = None,
        batch_size : int = Body(...,
                               description="生成图片数量",
                               example=""),
        styles : str = Body(...,
                               description="样式预设",
                               example=""),
):
    try:
        if width>=1024:
            width=1024
        if height>=1024:
            height=1024
        if batch_size>=2:
            batch_size=2
        if negative_prompt is None:
            negative_prompt = ""

        url = 'http://10.211.25.28:3000/v1/chat/completions'

        headers = {
            "Content-Type": "application/json",
            "Authorization": "Bearer sk-9aKmaj8vpFGB9OlRCdFaA7Dc352a40158d9578Bc29A0Ac6e"
        }


        prompt_en = ""
        data = {
                    "model": "SparkDesk",
                    "messages":  [
                        {
                            "role": "user",
                            "content": "以下句子请翻译成英文，如果是英文就不用翻译了：" + prompt
                        }
                    ],

                    "temperature": 0.7,
                    "top_p": 1,
                    "max_tokens": 4096,
                    "stop": [
                        "string"
                    ],
                    "user": "pdmi"
                }

        response = requests.post(url, json=data, headers=headers)  # 使用 JSON 格式发送数据
        if response.status_code == 200:
            # 解析响应内容
            response_data = response.json()  # 如果服务器返回 JSON 数据
            prompt_en = response_data['choices'][0]['message']['content']
            print(prompt_en)
        else:
            return BaseResponse(code=404, msg="模型翻译失败")


        negative_prompt_en = ""
        if len(negative_prompt) > 0:
            data_n = {
                "model": "SparkDesk",
                "messages": [
                    {
                        "role": "user",
                        "content": "以下句子请翻译成英文，如果是英文就不用翻译了：" + negative_prompt
                    }
                ],

                "temperature": 0.7,
                "top_p": 1,
                "max_tokens": 4096,
                "stop": [
                    "string"
                ],
                "user": "pdmi"
            }

            response = requests.post(url, json=data, headers=headers)  # 使用 JSON 格式发送数据
            if response.status_code == 200:
                # 解析响应内容
                response_data = response.json()  # 如果服务器返回 JSON 数据
                negative_prompt_en = response_data['choices'][0]['message']['content']
                print(negative_prompt_en)
            else:
                return BaseResponse(code=404, msg="模型翻译失败")


        # if response.status_code and response_n.status_code == 200:
        #     # 获取响应内容
        #     result = response.json()
        #     result_n = response_n.json()
        #
        #     prompt_en = result['response']
        #     negative_prompt_en = result_n['response']
        # else:
        #     return BaseResponse(code=404, msg="翻译过程出错")

        url = 'http://dm.hubpd.com:7005/sdapi/v1/txt2img'
        stylesList = []
        stylesList.append(styles)

        data = {
            "prompt": prompt_en,
            "negative_prompt": negative_prompt_en,
            "styles": stylesList,
            "seed": -1,
            "batch_size": batch_size,
            "steps": 50,
            "cfg_scale": 7,
            "width": width,
            "height": height,
            "sampler_index": "DPM++ 3M SDE Karras",
            "refiner_checkpoint": "sd_xl_refiner_1.0_0.9vae.safetensors"
        }

        logging.info(data)

        response = requests.post(url, headers=headers, json=data)

        if response.status_code == 200:
            # 获取响应内容
            result = response.json()
            images = result['images']
            return ListDocsResponse(data = images)
        else:
            return BaseResponse(code=404, msg="文生图接口访问失败")

    except Exception as e:
        return BaseResponse(code=404,msg=e)





def api_start(host, port, **kwargs):
    global app
    global local_doc_qa

    llm_model_ins = shared.loaderLLM()

    app = FastAPI()
    # Add CORS middleware to allow all origins
    # 在config.py中设置OPEN_DOMAIN=True，允许跨域
    # set OPEN_DOMAIN=True in config.py to allow cross-domain
    if OPEN_CROSS_DOMAIN:
        app.add_middleware(
            CORSMiddleware,
            allow_origins=["*"],
            allow_credentials=True,
            allow_methods=["*"],
            allow_headers=["*"],
        )
    # 修改了stream_chat的接口，直接通过ws://localhost:7861/local_doc_qa/stream_chat建立连接，在请求体中选择knowledge_base_id
    app.websocket("/local_doc_qa/stream_chat")(stream_chat)

    app.get("/", response_model=BaseResponse, summary="swagger 文档")(document)

    # 增加基于bing搜索的流式问答
    # 需要说明的是，如果想测试websocket的流式问答，需要使用支持websocket的测试工具，如postman,insomnia
    # 强烈推荐开源的insomnia
    # 在测试时选择new websocket request,并将url的协议改为ws,如ws://localhost:7861/local_doc_qa/stream_chat_bing
    app.websocket("/local_doc_qa/stream_chat_bing")(stream_chat_bing)
    app.post("/chat", response_model=ChatMessage, summary="与模型对话")(chat)
    app.post("/local_doc_qa/upload_file", response_model=BaseResponse, summary="上传文件到知识库")(upload_file)
    app.post("/local_doc_qa/upload_files", response_model=BaseResponse, summary="批量上传文件到知识库")(upload_files)
    app.post("/local_doc_qa/local_doc_chat", response_model=ChatMessage, summary="与知识库对话")(local_doc_chat)
    app.post("/local_doc_qa/bing_search_chat", response_model=ChatMessage, summary="与必应搜索对话")(bing_search_chat)
    app.get("/local_doc_qa/list_knowledge_base", response_model=ListDocsResponse, summary="获取知识库列表")(list_kbs)
    app.get("/local_doc_qa/list_files", response_model=ListDocsCreateResponse, summary="获取知识库内的文件列表")(list_docs)
    app.delete("/local_doc_qa/delete_knowledge_base", response_model=BaseResponse, summary="删除知识库")(delete_kb)
    app.delete("/local_doc_qa/delete_file", response_model=BaseResponse, summary="删除知识库内的文件")(delete_doc)
    app.post("/local_doc_qa/update_file", response_model=BaseResponse, summary="上传文件到知识库，并删除另一个文件")(update_doc)

    # 添加三个接口  登陆，返回注册信息，注销
    app.post("/local_doc_qa/login", response_model=BaseResponse, summary="用户登录")(login)
    app.post("/user/get_user_detail", response_model=UserDetailResponse, summary="返回用户详情接口")(user_detail)
    app.post("/user/logout", response_model=BaseResponse, summary="用户登出")(logout)

    # # 添加数据集接口
    app.post("/local_doc_qa/get_dataset_file_detail", response_model=DBMessage, summary="数据文件列表接口传入")(get_preview_urls)
    app.post("/local_doc_qa/get_dataset_data", response_model=DBMessage, summary="数据集检索接口")(get_dataset_data)
    app.post("/local_doc_qa/insert_data_value", response_model=BaseResponse, summary="加入价值观数据表接口")(insert_data_value)
    app.post("/local_doc_qa/insert_data_novalues", response_model=BaseResponse, summary="加入异常数据表接口")(insert_data_novalues)
    app.post("/local_doc_qa/insert_dataset", response_model=UIDResponse, summary="创建数据集接口")(insert_dataset)
    app.post("/local_doc_qa/search_data_values", response_model=DBMessage, summary="治理数据预览接口")(search_data_values)
    app.post("/local_doc_qa/search_data_novalues", response_model=DBMessage, summary="异常数据预览接口")(search_data_novalues)
    app.post("/local_doc_qa/get_dataset_data2", response_model=DBMessage, summary="数据集检索接口")(get_dataset_data2)
    app.post("/local_doc_qa/copyto_data_values", response_model=BaseResponse, summary="复制数据集加入治理")(copyto_data_values)
    app.post("/local_doc_qa/edit_data_values", response_model=BaseResponse, summary="编辑数据治理集")(edit_data_values)
    app.post("/local_doc_qa/edit_data_novalues", response_model=BaseResponse, summary="编辑异常数据集")(edit_data_novalues)
    app.post("/local_doc_qa/upload_dataset_files", response_model=BaseResponse, summary="数据集文件上传")(upload_dataset_files)
    app.post("/local_doc_qa/get_dataset_files", response_model=ListDocsResponse, summary="数据集文件下载")(get_dataset_files)
    app.post("/local_doc_qa/upload_dataset_files_temp", response_model=BaseResponse, summary="数据集文件上传暂存区")(upload_dataset_files_temp)

    #文生图
    app.post("/local_doc_qa/text2images", response_model=ListDocsResponse, summary="文生图")(text2images)
    # 知识库文件统计
    app.post("/local_doc_qa/knowleage_detail_count", response_model=UserDetailResponse, summary="知识库文件统计")(count_files_in_folder)











    # 访问文件
    # app.get("/local_doc_qa/get_file",  summary="访问文件")(get_file)
    # 与chat对话先访问本地chat知识库再访问大模型
    app.post("/chat_v2", response_model=ChatMessage, summary="与知识库和模型对话")(chat_v2)

    local_doc_qa = LocalDocQA()
    local_doc_qa.init_cfg(
        llm_model=llm_model_ins,
        embedding_model=EMBEDDING_MODEL,
        embedding_device=EMBEDDING_DEVICE,
        top_k=VECTOR_SEARCH_TOP_K,
    )
    if kwargs.get("ssl_keyfile") and kwargs.get("ssl_certfile"):
        uvicorn.run(app, workers=kwargs.get("worker_num"), host=host, port=port, ssl_keyfile=kwargs.get("ssl_keyfile"),
                    ssl_certfile=kwargs.get("ssl_certfile"))
    else:
        uvicorn.run(app, workers=kwargs.get("worker_num"), host=host, port=port)


if __name__ == "__main__":
    # 获取当前文件的绝对路径
    current_file_path = os.path.abspath(__file__)
    parent_directory_path = os.path.dirname(current_file_path)
    parent_directory = parent_directory_path.split('/')[-1]

    directory_parts = parent_directory.split('-')
    num_parts = len(directory_parts)
    if num_parts == 2:
        port_offset = 1
    else:
        port_offset = int(directory_parts[2])

    port = 7000 + port_offset

    parser.add_argument("--worker_num", type=int, default=1)
    parser.add_argument("--host", type=str, default="0.0.0.0")
    parser.add_argument("--port", type=int, default=port)
    parser.add_argument("--ssl_keyfile", type=str)
    parser.add_argument("--ssl_certfile", type=str)
    # 初始化消息

    args = parser.parse_args()
    args_dict = vars(args)
    shared.loaderCheckPoint = LoaderCheckPoint(args_dict)
    api_start(args.host, args.port, ssl_keyfile=args.ssl_keyfile, ssl_certfile=args.ssl_certfile)
